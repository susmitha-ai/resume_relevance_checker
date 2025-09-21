from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_enhanced_presentation():
    # Create presentation object
    prs = Presentation()
    
    # Set slide size to widescreen
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Enhanced Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background color
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(41, 128, 185)  # Blue background
    
    # Add title
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.33), Inches(2))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_p = title_frame.paragraphs[0]
    title_p.text = "AI-Powered Resume Relevance Checker"
    title_p.font.size = Pt(48)
    title_p.font.color.rgb = RGBColor(255, 255, 255)  # White text
    title_p.font.bold = True
    title_p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11.33), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.clear()
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = "Automated Recruitment Solution for Innomatics Research Labs"
    subtitle_p.font.size = Pt(24)
    subtitle_p.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    # Add demo link
    demo_box = slide1.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.33), Inches(1))
    demo_frame = demo_box.text_frame
    demo_frame.clear()
    demo_p = demo_frame.paragraphs[0]
    demo_p.text = "🚀 Live Demo: https://automatic-resume-relevance-checker.streamlit.app"
    demo_p.font.size = Pt(18)
    demo_p.font.color.rgb = RGBColor(255, 255, 255)
    demo_p.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Problem Statement with better design
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add background
    background = slide2.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(231, 76, 60)  # Red background
    
    # Add title
    title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_p = title_frame.paragraphs[0]
    title_p.text = "❌ The Challenge"
    title_p.font.size = Pt(36)
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.font.bold = True
    title_p.alignment = PP_ALIGN.CENTER
    
    # Add content with bullet points
    content_box = slide2.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(4.5))
    content_frame = content_box.text_frame
    content_frame.clear()
    
    problems = [
        "⏰ Time-consuming and inconsistent process",
        "📊 18-20 job requirements weekly with thousands of applications",
        "⏳ Delays in shortlisting candidates",
        "🔄 Inconsistent judgments across evaluators",
        "😰 High workload for placement staff",
        "📉 Reduced focus on interview preparation"
    ]
    
    for i, problem in enumerate(problems):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        p.text = problem
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_after = Pt(12)
    
    # Add impact statement
    impact_box = slide2.shapes.add_textbox(Inches(1), Inches(6), Inches(11.33), Inches(1))
    impact_frame = impact_box.text_frame
    impact_frame.clear()
    impact_p = impact_frame.paragraphs[0]
    impact_p.text = "💥 Impact: 80% of time spent on manual evaluation"
    impact_p.font.size = Pt(24)
    impact_p.font.color.rgb = RGBColor(255, 255, 255)
    impact_p.font.bold = True
    impact_p.alignment = PP_ALIGN.CENTER
    
    # Slide 3: Solution with green background
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add background
    background = slide3.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(39, 174, 96)  # Green background
    
    # Add title
    title_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_p = title_frame.paragraphs[0]
    title_p.text = "✅ Our Solution"
    title_p.font.size = Pt(36)
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.font.bold = True
    title_p.alignment = PP_ALIGN.CENTER
    
    # Add solution content
    content_box = slide3.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(4.5))
    content_frame = content_box.text_frame
    content_frame.clear()
    
    solutions = [
        "🤖 AI-Powered Resume Evaluation Platform",
        "",
        "🔍 Hybrid Scoring Algorithm",
        "   • Hard Match (60%) + Soft Match (40%)",
        "   • Keyword matching + Semantic analysis",
        "",
        "📊 Multi-Modal Analysis",
        "   • 5 different analysis perspectives",
        "   • Standard, ATS, Performance, Strength, Comparison",
        "",
        "⚡ Real-Time Processing",
        "   • 2-3 seconds per resume",
        "   • Batch processing capabilities"
    ]
    
    for i, solution in enumerate(solutions):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        p.text = solution
        if solution.startswith("🤖") or solution.startswith("🔍") or solution.startswith("📊") or solution.startswith("⚡"):
            p.font.size = Pt(22)
            p.font.bold = True
        else:
            p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_after = Pt(8)
    
    # Add benefits
    benefits_box = slide3.shapes.add_textbox(Inches(1), Inches(6), Inches(11.33), Inches(1))
    benefits_frame = benefits_box.text_frame
    benefits_frame.clear()
    benefits_p = benefits_frame.paragraphs[0]
    benefits_p.text = "🎯 Key Benefits: 80% time reduction • Consistent scoring • Personalized feedback"
    benefits_p.font.size = Pt(20)
    benefits_p.font.color.rgb = RGBColor(255, 255, 255)
    benefits_p.font.bold = True
    benefits_p.alignment = PP_ALIGN.CENTER
    
    # Slide 4: Features with purple background
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add background
    background = slide4.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(155, 89, 182)  # Purple background
    
    # Add title
    title_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_p = title_frame.paragraphs[0]
    title_p.text = "🚀 Core Features"
    title_p.font.size = Pt(36)
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.font.bold = True
    title_p.alignment = PP_ALIGN.CENTER
    
    # Add features in two columns
    # Left column
    left_box = slide4.shapes.add_textbox(Inches(0.5), Inches(2), Inches(6), Inches(4.5))
    left_frame = left_box.text_frame
    left_frame.clear()
    
    left_features = [
        "📊 Automated Analysis:",
        "✅ Relevance Scoring (0-100)",
        "✅ Gap Analysis",
        "✅ Fit Verdict (High/Medium/Low)",
        "✅ Personalized Feedback",
        "",
        "👥 User Management:",
        "✅ Secure Authentication",
        "✅ Demo Accounts",
        "✅ Analysis History",
        "✅ Data Export (CSV/JSON)"
    ]
    
    for i, feature in enumerate(left_features):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = feature
        if feature.endswith(":"):
            p.font.size = Pt(20)
            p.font.bold = True
        else:
            p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_after = Pt(6)
    
    # Right column
    right_box = slide4.shapes.add_textbox(Inches(6.5), Inches(2), Inches(6), Inches(4.5))
    right_frame = right_box.text_frame
    right_frame.clear()
    
    right_features = [
        "🔍 Advanced Analysis Modes:",
        "• Standard Analysis",
        "• ATS Score Analysis",
        "• Performance Prediction",
        "• Strength Analysis",
        "• Comparison Dashboard",
        "",
        "🎯 Analysis Types:",
        "• Comprehensive matching",
        "• Applicant tracking compatibility",
        "• Interview likelihood",
        "• Multi-dimensional evaluation",
        "• Side-by-side ranking"
    ]
    
    for i, feature in enumerate(right_features):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = feature
        if feature.endswith(":"):
            p.font.size = Pt(20)
            p.font.bold = True
        else:
            p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_after = Pt(6)
    
    # Slide 5: Live Demo with orange background
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add background
    background = slide5.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(230, 126, 34)  # Orange background
    
    # Add title
    title_box = slide5.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_p = title_frame.paragraphs[0]
    title_p.text = "🎬 Live Demo"
    title_p.font.size = Pt(36)
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.font.bold = True
    title_p.alignment = PP_ALIGN.CENTER
    
    # Add demo steps
    steps_box = slide5.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(3))
    steps_frame = steps_box.text_frame
    steps_frame.clear()
    
    steps = [
        "1️⃣ Access: https://automatic-resume-relevance-checker.streamlit.app",
        "2️⃣ Login: Use demo account (demo_hr/demo123)",
        "3️⃣ Upload JD: Paste job description",
        "4️⃣ Upload Resumes: Multiple PDF/DOCX files",
        "5️⃣ Run Analysis: Select analysis mode",
        "6️⃣ View Results: Scores, verdicts, feedback"
    ]
    
    for i, step in enumerate(steps):
        if i == 0:
            p = steps_frame.paragraphs[0]
        else:
            p = steps_frame.add_paragraph()
        p.text = step
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_after = Pt(12)
    
    # Add demo accounts
    accounts_box = slide5.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.33), Inches(1.5))
    accounts_frame = accounts_box.text_frame
    accounts_frame.clear()
    accounts_p = accounts_frame.paragraphs[0]
    accounts_p.text = "🎭 Demo Accounts:\ndemo_hr/demo123 • demo_recruiter/demo123 • demo_manager/demo123"
    accounts_p.font.size = Pt(18)
    accounts_p.font.color.rgb = RGBColor(255, 255, 255)
    accounts_p.font.bold = True
    accounts_p.alignment = PP_ALIGN.CENTER
    
    # Slide 6: Business Impact with teal background
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add background
    background = slide6.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(26, 188, 156)  # Teal background
    
    # Add title
    title_box = slide6.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_p = title_frame.paragraphs[0]
    title_p.text = "💼 Business Impact"
    title_p.font.size = Pt(36)
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.font.bold = True
    title_p.alignment = PP_ALIGN.CENTER
    
    # Add impact metrics in grid
    metrics = [
        ("👥 For Placement Teams", "80% time savings\nStandardized scoring\nHandle increased volumes"),
        ("🎓 For Students", "Immediate feedback\nSkill gap identification\nCareer guidance"),
        ("🏢 For Organizations", "Reduced costs\nFaster hiring\nBetter matches"),
        ("📊 ROI Metrics", "60-70% cost reduction\n80% time savings\n95% satisfaction")
    ]
    
    for i, (title, content) in enumerate(metrics):
        x = Inches(0.5 + (i % 2) * 6.5)
        y = Inches(2 + (i // 2) * 2.5)
        
        # Add metric box
        metric_box = slide6.shapes.add_textbox(x, y, Inches(6), Inches(2))
        metric_frame = metric_box.text_frame
        metric_frame.clear()
        
        # Add title
        title_p = metric_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(18)
        title_p.font.color.rgb = RGBColor(255, 255, 255)
        title_p.font.bold = True
        title_p.alignment = PP_ALIGN.CENTER
        
        # Add content
        content_p = metric_frame.add_paragraph()
        content_p.text = content
        content_p.font.size = Pt(14)
        content_p.font.color.rgb = RGBColor(255, 255, 255)
        content_p.alignment = PP_ALIGN.CENTER
    
    # Slide 7: Thank You with gradient-like effect
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add background
    background = slide7.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(52, 73, 94)  # Dark blue-gray background
    
    # Add title
    title_box = slide7.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(2))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_p = title_frame.paragraphs[0]
    title_p.text = "🙏 Thank You!"
    title_p.font.size = Pt(48)
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.font.bold = True
    title_p.alignment = PP_ALIGN.CENTER
    
    # Add key takeaways
    takeaways_box = slide7.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.33), Inches(2))
    takeaways_frame = takeaways_box.text_frame
    takeaways_frame.clear()
    takeaways_p = takeaways_frame.paragraphs[0]
    takeaways_p.text = "✨ Innovative AI-powered solution • Production-ready deployment • Measurable business impact"
    takeaways_p.font.size = Pt(20)
    takeaways_p.font.color.rgb = RGBColor(255, 255, 255)
    takeaways_p.alignment = PP_ALIGN.CENTER
    
    # Add contact info
    contact_box = slide7.shapes.add_textbox(Inches(1), Inches(6), Inches(11.33), Inches(1))
    contact_frame = contact_box.text_frame
    contact_frame.clear()
    contact_p = contact_frame.paragraphs[0]
    contact_p.text = "🚀 Live Demo: https://automatic-resume-relevance-checker.streamlit.app"
    contact_p.font.size = Pt(16)
    contact_p.font.color.rgb = RGBColor(255, 255, 255)
    contact_p.alignment = PP_ALIGN.CENTER
    
    # Save the enhanced presentation
    prs.save('Enhanced_Resume_Relevance_Checker_Presentation.pptx')
    print("✅ Enhanced PowerPoint presentation created successfully!")
    print("📁 File saved as: Enhanced_Resume_Relevance_Checker_Presentation.pptx")
    print("🎨 Features: Colorful backgrounds, better typography, visual appeal")
    print("📥 Ready for download!")

if __name__ == "__main__":
    create_enhanced_presentation()
